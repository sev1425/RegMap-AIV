# ==========================================================
# REGMAP AI ENTERPRISE
# Universal Document Loader
# ==========================================================

from pathlib import Path
import fitz
from PIL import Image
import pytesseract
from docx import Document
from pptx import Presentation


class UniversalDocumentLoader:

    def __init__(self, file_path):

        self.file_path = Path(file_path)

    # ======================================================
    # Main Loader
    # ======================================================

    def load(self):

        if not self.file_path.exists():
            raise FileNotFoundError(
                f"File not found: {self.file_path}"
            )

        extension = self.file_path.suffix.lower()

        if extension == ".pdf":
            return self.read_pdf()

        elif extension == ".docx":
            return self.read_docx()

        elif extension == ".txt":
            return self.read_txt()

        elif extension in [".png", ".jpg", ".jpeg"]:
            return self.read_image()

        elif extension == ".pptx":
            return self.read_ppt()

        # Fallback for ANY other file type: Try to read as plain text
        try:
            return self.read_txt()
        except Exception as e:
            return f"[Error: Unable to parse unsupported binary format {extension}]"

    # ======================================================
    # PDF
    # ======================================================

    def read_pdf(self):

        print(f"[INFO] Reading PDF: {self.file_path}")

        text = ""

        with fitz.open(self.file_path) as pdf:

            print(f"[INFO] Pages: {pdf.page_count}")

            for index, page in enumerate(pdf):

                page_text = page.get_text()

                text += page_text

                print(
                    f"[INFO] Page {index+1}: "
                    f"{len(page_text)} characters"
                )

        print(f"[INFO] Total characters: {len(text)}")

        return text

    # ======================================================
    # DOCX
    # ======================================================

    def read_docx(self):

        text = ""

        doc = Document(self.file_path)

        for paragraph in doc.paragraphs:

            text += paragraph.text + "\n"

        return text

    # ======================================================
    # TXT
    # ======================================================

    def read_txt(self):

        with open(
            self.file_path,
            "r",
            encoding="utf-8"
        ) as file:

            return file.read()

    # ======================================================
    # IMAGE OCR
    # ======================================================

    def read_image(self):
        try:
            image = Image.open(self.file_path)
            text = pytesseract.image_to_string(image)
            if not text.strip():
                return "[OCR_WARNING] No readable text was found in this image."
            return text
        except Exception as e:
            print(f"[OCR Warning] Tesseract not found or failed: {e}")
            return (
                f"[OCR_UNAVAILABLE] Could not extract text from this image "
                f"(OCR engine unavailable or failed: {e}). "
                f"Please upload a text-based document instead, or ensure "
                f"Tesseract OCR is installed."
            )

    # ======================================================
    # PPTX
    # ======================================================

    def read_ppt(self):

        text = ""

        presentation = Presentation(self.file_path)

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text += shape.text + "\n"

        return text